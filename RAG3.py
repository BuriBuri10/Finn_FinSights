import os
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from langchain.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
from langchain.chains.qa_with_sources import load_qa_with_sources_chain
from langchain_groq import ChatGroq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import textwrap
import tempfile
import re

# Streamlit UI
st.set_page_config(page_title="AutoAnalyst")
st.title("📊 FinSights -- Finn </autoAnalyst>")

uploaded_files = st.file_uploader("Upload one or more financial reports (PDF)", type=["pdf"], accept_multiple_files=True)
query = st.text_input("Ask a financial question (e.g., 'Summarize revenue trends'):")

FEW_SHOT_EXAMPLES = """
You are a financial assistant. Based on the document context and user question, respond with structured and concise summaries.

Example 1:
Question: Summarize the revenue trends.
Answer:
- Revenue increased by 12% YoY, reaching $3.2B in Q4.
- Growth driven by strong performance in the APAC region.

Example 2:
Question: What are the key risks highlighted?
Answer:
- Exposure to volatile foreign exchange rates.
- Regulatory risks in European operations.
- Supply chain disruptions due to geopolitical tensions.
"""

custom_prompt = PromptTemplate(
    input_variables=["context", "question"],
    template=FEW_SHOT_EXAMPLES + "\nContext: {context}\nQuestion: {question}\nAnswer:"
)

def load_and_split_pdfs(files):
    documents = []
    for file in files:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                tmp_file.write(file.read())
                loader = PyPDFLoader(tmp_file.name)
                docs = loader.load()
                documents.extend(docs)
        except Exception as e:
            st.warning(f"⚠️ Failed to process file {file.name}: {e}")
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    return splitter.split_documents(documents)

def build_vector_store(docs):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_documents(docs, embeddings)

def create_qa_chain(vector_store):
    retriever = vector_store.as_retriever()
    llm = ChatGroq(model_name="deepseek-r1-distill-llama-70b", api_key="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")
    return RetrievalQA.from_chain_type(llm=llm, retriever=retriever, chain_type_kwargs={"prompt": custom_prompt})

def extract_financial_metrics(text):
    metrics = {}
    patterns = {
        "Revenue": r"revenue[^\d]*(\d+[,.]?\d*[\w]*)",
        "Profit": r"profit[^\d]*(\d+[,.]?\d*[\w]*)",
        "Net Income": r"net income[^\d]*(\d+[,.]?\d*[\w]*)",
        "EBITDA": r"EBITDA[^\d]*(\d+[,.]?\d*[\w]*)",
        "Expenses": r"expenses[^\d]*(\d+[,.]?\d*[\w]*)",
        "Cash Flow": r"cash flow[^\d]*(\d+[,.]?\d*[\w]*)",
    }
    for key, pattern in patterns.items():
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            metrics[key] = match.group(1)
    return metrics

def plot_metrics(metrics):
    if not metrics:
        return
    df = pd.DataFrame(metrics.items(), columns=["Metric", "Value"])
    fig, ax = plt.subplots()
    ax.barh(df["Metric"], [float(re.sub(r'[^\d.]', '', v)) for v in df["Value"]])
    ax.set_xlabel("Amount")
    st.pyplot(fig)

def generate_pptx(title, content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]  # Fully blank layout

    # Brand styling
    bg_color = RGBColor(25, 50, 112)
    accent_color = RGBColor(91, 155, 213)
    font_color = RGBColor(255, 255, 255)

    # Title Slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "AutoAnalyst Financial Summary"
    title_slide.placeholders[1].text = f"Query: {title}"
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = bg_color

    # Clean sentence splitting
    sentences = re.split(r'(?<=[.!?])\s+', content.strip())
    grouped_sentences = [sentences[i:i + 5] for i in range(0, len(sentences), 5)]

    for index, chunk in enumerate(grouped_sentences):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.text = f"Insight {index + 1}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = font_color
        p.alignment = 1  # Center

        # Content Box
        left_margin = Inches(0.8)
        top_margin = Inches(1.2)
        width = Inches(8)
        height = Inches(4.5)
        content_box = slide.shapes.add_textbox(left_margin, top_margin, width, height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Determine font size based on number of lines
        num_lines = len(chunk)
        if num_lines <= 3:
            font_size = Pt(18)
        elif num_lines <= 5:
            font_size = Pt(16)
        else:
            font_size = Pt(14)

        for sentence in chunk:
            p = content_frame.add_paragraph()
            p.text = sentence
            p.level = 0
            p.font.size = font_size
            p.font.color.rgb = font_color
            p.space_after = Pt(5)

        # Footer Tag
        tag_left = Inches(7)
        tag_top = Inches(6.8)
        tag_width = Inches(2)
        tag_height = Inches(0.4)
        tag = slide.shapes.add_textbox(tag_left, tag_top, tag_width, tag_height)
        tag_frame = tag.text_frame
        tag_p = tag_frame.paragraphs[0]
        tag_p.text = "AutoAnalyst"
        tag_p.font.size = Pt(11)
        tag_p.font.bold = True
        tag_p.font.color.rgb = font_color

    # Financial Metrics Chart Slide
    metrics = extract_financial_metrics(content)
    if metrics:
        chart_slide = prs.slides.add_slide(blank_slide_layout)
        chart_slide.background.fill.solid()
        chart_slide.background.fill.fore_color.rgb = bg_color

        chart_title = chart_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        chart_title_tf = chart_title.text_frame
        chart_title_tf.text = "Key Financial Metrics"
        chart_title_tf.paragraphs[0].font.size = Pt(20)
        chart_title_tf.paragraphs[0].font.color.rgb = font_color
        chart_title_tf.paragraphs[0].alignment = 1

        chart_data = CategoryChartData()
        chart_data.categories = list(metrics.keys())
        values = [float(re.sub(r'[^\d.]', '', v)) for v in metrics.values()]
        chart_data.add_series('Values', values)

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        chart_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)

    tmp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_pptx.name)
    return tmp_pptx.name

if uploaded_files and query:
    with st.spinner("Processing documents and generating insights..."):
        documents = load_and_split_pdfs(uploaded_files)
        vector_store = build_vector_store(documents)
        qa_chain = create_qa_chain(vector_store)

        response = qa_chain.run(query)
        st.subheader("📌 Answer:")
        st.write(response)

        metrics = extract_financial_metrics(response)
        if metrics:
            st.markdown("### 📈 Key Financial Metrics Identified:")
            st.write(metrics)
            plot_metrics(metrics)

        pptx_path = generate_pptx(query, response)
        with open(pptx_path, "rb") as f:
            st.download_button("Download Summary as PowerPoint", f, file_name="financial_summary.pptx")

        st.markdown("---")
        st.markdown("### 📊 Financial Insights Dashboard")
        if metrics:
            st.dataframe(pd.DataFrame(metrics.items(), columns=["Metric", "Extracted Value"]))
